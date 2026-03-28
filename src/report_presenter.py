"""Meridant Insight — PowerPoint executive readout export.

Generates a consulting-grade .pptx deck from a completed assessment.
Returns bytes suitable for st.download_button().

Slides:
  1. Cover — client name, engagement, consultant, date, Meridant Insight branding
  2. Executive Summary — narrative + key stats
  3. Maturity Heatmap — domain × L1–L5 table
  4. Top Gaps — P1 capability gaps table
  5. Transformation Roadmap — phase timeline table (or placeholder)
  6. Next Steps — editable placeholder bullets
"""

from __future__ import annotations

import io
from datetime import date
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ── Brand colours ─────────────────────────────────────────────────────────────
_NAVY       = RGBColor(0x0F, 0x27, 0x44)
_ACCENT     = RGBColor(0x25, 0x63, 0xEB)
_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
_GRAY_MID   = RGBColor(0x37, 0x41, 0x51)
_GRAY_LIGHT = RGBColor(0x6B, 0x72, 0x80)
_GREEN_BG   = RGBColor(0xCC, 0xFB, 0xF1)
_AMBER_BG   = RGBColor(0xFE, 0xF3, 0xC7)
_RED_TXT    = RGBColor(0xDC, 0x26, 0x26)
_AMBER_TXT  = RGBColor(0xD9, 0x77, 0x06)
_GREEN_TXT  = RGBColor(0x16, 0xA3, 0x4A)

# Slide dimensions (widescreen 16:9)
_W = Inches(13.33)
_H = Inches(7.5)

# Safe area margins
_ML = Inches(0.6)
_MT = Inches(0.55)
_CW = _W - 2 * _ML
_CH = _H - 2 * _MT


def _prs() -> Presentation:
    """Create a blank widescreen presentation."""
    prs = Presentation()
    prs.slide_width  = _W
    prs.slide_height = _H
    return prs


def _blank_slide(prs: Presentation):
    """Add a blank-layout slide."""
    blank = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(blank)


def _add_textbox(slide, left, top, width, height, text: str,
                 size: int = 18, bold: bool = False,
                 color: RGBColor = None, align=PP_ALIGN.LEFT,
                 wrap: bool = True) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _add_rect(slide, left, top, width, height, fill: RGBColor):
    """Add a filled rectangle (no border)."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def _level_score(avg_score, level: int) -> float:
    if avg_score is None:
        return 0.0
    return max(0.0, min(1.0, float(avg_score) - (level - 1)))


def _cell_color(score: float) -> RGBColor:
    if score >= 1.0:
        return _GREEN_BG
    if score > 0.0:
        return _AMBER_BG
    return RGBColor(0xFF, 0xFF, 0xFF)


# ── Slide builders ─────────────────────────────────────────────────────────────

def _slide_cover(prs: Presentation, client_name: str, engagement_name: str,
                 use_case_name: str, consultant_name: str) -> None:
    slide = _blank_slide(prs)

    # Navy background panel (left half)
    _add_rect(slide, 0, 0, Inches(5.5), _H, _NAVY)

    # "meridant insight" label
    _add_textbox(slide, _ML, Inches(1.5), Inches(4.5), Inches(0.5),
                 "meridant insight", size=11, bold=True, color=_ACCENT)

    # Client name (large)
    _add_textbox(slide, _ML, Inches(2.1), Inches(4.5), Inches(1.8),
                 client_name, size=32, bold=True, color=_WHITE, wrap=True)

    # Engagement
    _add_textbox(slide, _ML, Inches(3.9), Inches(4.5), Inches(0.5),
                 engagement_name, size=14, color=RGBColor(0xB0, 0xBE, 0xC5))

    # Use case
    _add_textbox(slide, _ML, Inches(4.5), Inches(4.5), Inches(0.4),
                 f"Use Case: {use_case_name}", size=10, color=_GRAY_LIGHT)

    # Consultant + date
    footer_text = f"{consultant_name}  ·  {date.today().strftime('%d %B %Y')}" if consultant_name \
                  else date.today().strftime("%d %B %Y")
    _add_textbox(slide, _ML, Inches(6.5), Inches(4.5), Inches(0.4),
                 footer_text, size=9, color=_GRAY_LIGHT)


def _slide_exec_summary(prs: Presentation, client_name: str,
                         findings_narrative: str, dom_scores: list,
                         cap_findings: list) -> None:
    slide = _blank_slide(prs)

    # Header bar
    _add_rect(slide, 0, 0, _W, Inches(0.8), _NAVY)
    _add_textbox(slide, _ML, Inches(0.1), Inches(8), Inches(0.6),
                 "Executive Summary", size=18, bold=True, color=_WHITE)

    # Stats row
    total_domains = len(dom_scores)
    total_caps    = len(cap_findings)
    scored_caps   = [c for c in cap_findings if c.get("avg_score") is not None]
    overall       = round(sum(c["avg_score"] for c in scored_caps) / len(scored_caps), 1) \
                    if scored_caps else None

    stats = [
        (str(total_domains), "Domains Assessed"),
        (str(total_caps),    "Capabilities"),
        (f"{overall:.1f}" if overall else "—", "Overall Score"),
    ]
    for i, (val, lbl) in enumerate(stats):
        x = _ML + Inches(i * 2.2)
        _add_textbox(slide, x, Inches(0.9), Inches(2.0), Inches(0.55),
                     val, size=26, bold=True, color=_NAVY)
        _add_textbox(slide, x, Inches(1.45), Inches(2.0), Inches(0.35),
                     lbl, size=9, color=_GRAY_LIGHT)

    # Narrative (truncated to ~500 chars on slide; full text in speaker notes)
    if findings_narrative:
        truncated = findings_narrative[:500].rsplit(" ", 1)[0] + "…" \
                    if len(findings_narrative) > 500 else findings_narrative
        _add_textbox(slide, _ML, Inches(1.9), _CW, Inches(4.5),
                     truncated, size=11, color=_GRAY_MID, wrap=True)
        # Full text in notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = findings_narrative
    else:
        _add_textbox(slide, _ML, Inches(1.9), _CW, Inches(1),
                     "No executive summary available.", size=11, color=_GRAY_LIGHT)


def _slide_heatmap(prs: Presentation, dom_scores: list) -> None:
    """Domain × maturity level heatmap table slide."""
    slide = _blank_slide(prs)

    _add_rect(slide, 0, 0, _W, Inches(0.8), _NAVY)
    _add_textbox(slide, _ML, Inches(0.1), Inches(8), Inches(0.6),
                 "Maturity Heatmap", size=18, bold=True, color=_WHITE)

    if not dom_scores:
        _add_textbox(slide, _ML, Inches(1.5), _CW, Inches(1),
                     "No domain findings available.", size=12, color=_GRAY_LIGHT)
        return

    levels = [1, 2, 3, 4, 5]
    level_names = ["Ad Hoc", "Defined", "Integrated", "Intelligent", "Adaptive"]

    rows = len(dom_scores) + 1  # +1 for header
    cols = len(levels) + 1      # +1 for domain name

    col_w = Inches(1.9)
    row_h = Inches(0.45)
    tbl_w = col_w * cols
    tbl_x = (_W - tbl_w) / 2
    tbl_y = Inches(1.0)

    table = slide.shapes.add_table(rows, cols, tbl_x, tbl_y, tbl_w, row_h * rows).table

    # Header row
    hdr_row = table.rows[0]
    hdr_row.cells[0].text = "Domain"
    for i, (lv, nm) in enumerate(zip(levels, level_names)):
        hdr_row.cells[i + 1].text = f"L{lv} {nm}"

    for cell in hdr_row.cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = _NAVY
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = Pt(9)
                run.font.bold = True
                run.font.color.rgb = _WHITE

    # Data rows
    for ri, d in enumerate(dom_scores):
        row = table.rows[ri + 1]
        row.cells[0].text = d.get("domain", "")
        for ci, lv in enumerate(levels):
            cell = row.cells[ci + 1]
            score = _level_score(d.get("avg_score"), lv)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _cell_color(score)
            cell_text = f"{score:.0%}" if score > 0 else ""
            cell.text = cell_text
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(8)
                    run.font.color.rgb = _GRAY_MID

        # Domain cell style
        row.cells[0].fill.solid()
        row.cells[0].fill.fore_color.rgb = RGBColor(0xF9, 0xFA, 0xFB)
        for para in row.cells[0].text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(9)
                run.font.bold = True
                run.font.color.rgb = _NAVY


def _slide_top_gaps(prs: Presentation, recommendations: list) -> None:
    """Top P1 capability gaps table slide."""
    slide = _blank_slide(prs)

    _add_rect(slide, 0, 0, _W, Inches(0.8), _NAVY)
    _add_textbox(slide, _ML, Inches(0.1), Inches(8), Inches(0.6),
                 "Priority Gaps — P1 Foundation", size=18, bold=True, color=_WHITE)

    p1_recs = [r for r in recommendations if r.get("priority_tier") == "P1"][:10]
    if not p1_recs:
        _add_textbox(slide, _ML, Inches(1.5), _CW, Inches(1),
                     "No P1 recommendations available.", size=12, color=_GRAY_LIGHT)
        return

    rows = len(p1_recs) + 1
    cols = 5
    col_widths_in = [3.5, 2.2, 1.2, 1.2, 1.0]
    tbl_w = sum(Inches(w) for w in col_widths_in)
    tbl_x = (_W - tbl_w) / 2
    tbl_y = Inches(1.0)
    row_h = Inches(0.42)

    table = slide.shapes.add_table(rows, cols, tbl_x, tbl_y, tbl_w, row_h * rows).table

    # Set column widths
    for ci, w in enumerate(col_widths_in):
        for ri in range(rows):
            table.rows[ri].cells[ci].width = Inches(w)

    headers = ["Capability", "Domain", "Score", "Target", "Gap"]
    hdr_row = table.rows[0]
    for ci, h in enumerate(headers):
        hdr_row.cells[ci].text = h
        hdr_row.cells[ci].fill.solid()
        hdr_row.cells[ci].fill.fore_color.rgb = _NAVY
        for para in hdr_row.cells[ci].text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(9)
                run.font.bold = True
                run.font.color.rgb = _WHITE

    for ri, rec in enumerate(p1_recs):
        row = table.rows[ri + 1]
        row.cells[0].text = rec.get("capability_name", "")
        row.cells[1].text = rec.get("domain", "")
        cur = rec.get("current_score")
        row.cells[2].text = f"{cur:.1f}" if cur is not None else "—"
        row.cells[3].text = str(rec.get("target_maturity") or "—")
        gap = rec.get("gap") or 0
        row.cells[4].text = f"{gap:.1f}"
        for ci in range(cols):
            row.cells[ci].fill.solid()
            row.cells[ci].fill.fore_color.rgb = RGBColor(0xF9, 0xFA, 0xFB) if ri % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            for para in row.cells[ci].text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(9)
                    run.font.color.rgb = _GRAY_MID
        # Gap cell red if large
        if gap >= 1.0:
            for para in row.cells[4].text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = _RED_TXT
                    run.font.bold = True


def _slide_roadmap(prs: Presentation, roadmap: dict | None) -> None:
    """Transformation roadmap phase timeline slide."""
    slide = _blank_slide(prs)

    _add_rect(slide, 0, 0, _W, Inches(0.8), _NAVY)
    _add_textbox(slide, _ML, Inches(0.1), Inches(8), Inches(0.6),
                 "Transformation Roadmap", size=18, bold=True, color=_WHITE)

    if not roadmap or not roadmap.get("phases"):
        _add_textbox(slide, _ML, Inches(1.5), _CW, Inches(1.5),
                     "Roadmap not generated.\n\nGenerate a roadmap in Step 6 of the Assessment wizard, then re-export.",
                     size=12, color=_GRAY_LIGHT, wrap=True)
        return

    phases = roadmap["phases"]
    rows = len(phases) + 1
    cols = 4
    tbl_w = _CW
    tbl_x = _ML
    tbl_y = Inches(1.0)
    row_h = Inches(0.8)

    table = slide.shapes.add_table(rows, cols, tbl_x, tbl_y, tbl_w, row_h * rows).table
    col_w_in = [1.0, 2.5, 5.5, 2.0]
    for ci, w in enumerate(col_w_in):
        for ri in range(rows):
            table.rows[ri].cells[ci].width = Inches(w)

    headers = ["Phase", "Name", "Key Initiatives", "Timeline"]
    hdr_row = table.rows[0]
    for ci, h in enumerate(headers):
        hdr_row.cells[ci].text = h
        hdr_row.cells[ci].fill.solid()
        hdr_row.cells[ci].fill.fore_color.rgb = _NAVY
        for para in hdr_row.cells[ci].text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(9)
                run.font.bold = True
                run.font.color.rgb = _WHITE

    for ri, phase in enumerate(phases):
        row = table.rows[ri + 1]
        initiatives = phase.get("initiatives", [])
        init_names  = ", ".join(i["name"] for i in initiatives[:3])
        if len(initiatives) > 3:
            init_names += f" +{len(initiatives) - 3} more"
        timeline_txt = f"W{phase.get('start_week','?')}–W{phase.get('end_week','?')}"
        row.cells[0].text = phase.get("id", "")
        row.cells[1].text = phase.get("name", "")
        row.cells[2].text = init_names or phase.get("description", "")
        row.cells[3].text = timeline_txt
        bg = RGBColor(0xF9, 0xFA, 0xFB) if ri % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
        for ci in range(cols):
            row.cells[ci].fill.solid()
            row.cells[ci].fill.fore_color.rgb = bg
            for para in row.cells[ci].text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(9)
                    run.font.color.rgb = _GRAY_MID
        for para in row.cells[0].text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.color.rgb = _NAVY


def _slide_next_steps(prs: Presentation, client_name: str) -> None:
    slide = _blank_slide(prs)

    _add_rect(slide, 0, 0, _W, Inches(0.8), _NAVY)
    _add_textbox(slide, _ML, Inches(0.1), Inches(8), Inches(0.6),
                 "Next Steps", size=18, bold=True, color=_WHITE)

    bullets = [
        "Review P1 recommendations and confirm capability ownership",
        "Confirm resourcing and prioritisation of foundation initiatives",
        f"Schedule maturity re-assessment for {client_name} in 6–12 months",
        "[Consultant: customise these steps for the client engagement]",
    ]
    _add_textbox(slide, _ML, Inches(1.0), _CW, Inches(4.0),
                 "\n\n".join(f"•  {b}" for b in bullets),
                 size=14, color=_GRAY_MID, wrap=True)

    # Footer
    _add_textbox(slide, _ML, Inches(6.8), _CW, Inches(0.4),
                 "Meridant Insight  ·  Confidential",
                 size=9, color=_GRAY_LIGHT, align=PP_ALIGN.CENTER)


# ── Public API ─────────────────────────────────────────────────────────────────

def generate_pptx_report(
    client_name: str,
    engagement_name: str,
    use_case_name: str,
    consultant_name: str,
    findings_narrative: str,
    dom_scores: list[dict],
    cap_findings: list[dict],
    recommendations: list[dict],
    roadmap: dict | None = None,
) -> bytes:
    """
    Generate a PowerPoint executive readout.

    Returns bytes suitable for st.download_button(
        mime='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    ).
    """
    prs = _prs()

    _slide_cover(prs, client_name, engagement_name, use_case_name, consultant_name)
    _slide_exec_summary(prs, client_name, findings_narrative, dom_scores, cap_findings)
    _slide_heatmap(prs, dom_scores)
    _slide_top_gaps(prs, recommendations or [])
    _slide_roadmap(prs, roadmap)
    _slide_next_steps(prs, client_name)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
